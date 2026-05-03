"""Tools the LLM agents can call.

Every tool is a read-only lookup against the workbench's existing in-process
state — workspace data, datasets, models, runs, plots, the workflow validator.
The OpenAI / agents SDK calls `handle_tool_call(name, args)` which dispatches
to the implementation here.

Filters / writes (`apply_tile_filter`) intentionally exist so the Tile Tuner can
materially change the tile when the analyst clicks an action chip.
"""
from __future__ import annotations

import contextvars
import json
import math
from typing import Any

import pandas as pd

from routers.datasets import _DATASETS, _read_dataframe, _resolve_path, _synthesize_sample
from routers.models_registry import _MODELS
from routers.plots import _PLOTS, _apply_filters
from routers.scenarios import _RUNS, _SCENARIOS
from services.workspace_data import get_workspace


# ── Per-request context (entity_kind / entity_id / function_id) ──────────
# The chat router pushes the request's bound entity into this contextvar
# before invoking the agent. Mutation tools fall back to it when the model
# omits `target_id` / `target_kind`, so a user who clicked "Tune" on a tile
# can ask "switch to bar" without the model needing to echo the tile id.
_REQUEST_CTX: contextvars.ContextVar[dict] = contextvars.ContextVar(
    "_cma_agent_request_ctx", default={},
)


def set_request_context(ctx: dict) -> contextvars.Token:
    """Set the active request context. Caller should reset() the token when
    the agent run finishes (use try/finally)."""
    return _REQUEST_CTX.set(dict(ctx or {}))


def reset_request_context(token: contextvars.Token) -> None:
    _REQUEST_CTX.reset(token)


def _ctx_target() -> tuple[str | None, str | None]:
    ctx = _REQUEST_CTX.get() or {}
    kind = ctx.get("entity_kind")
    eid = ctx.get("entity_id")
    if kind in ("tile", "analytic_def") and eid:
        return kind, eid
    return None, None


def _ctx_tile_id() -> str:
    kind, eid = _ctx_target()
    return eid if kind == "tile" else ""


def _ctx_id_for(expected_kind: str) -> str:
    """Generic single-id fallback for a specific entity kind. Returns the
    request's bound entity_id only when its kind matches. Used by lookup
    tools (get_dataset_preview, profile_dataset, get_model, …) so the
    model can omit / mistype the id and the tool still hits the right
    object that the analyst was looking at when they sent the message."""
    ctx = _REQUEST_CTX.get() or {}
    kind = ctx.get("entity_kind")
    eid = ctx.get("entity_id")
    return eid if kind == expected_kind and eid else ""


def _ctx_function_id() -> str:
    return (_REQUEST_CTX.get() or {}).get("function_id") or ""


# OpenAI-format tool schemas — these are advertised to the LLM
OPENAI_TOOLS: list[dict[str, Any]] = [
    {
        "type": "function",
        "function": {
            "name": "get_workspace",
            "description": "Fetch the live Overview workspace snapshot for a function — KPIs, charts, tables, insights.",
            "parameters": {
                "type": "object",
                "properties": {
                    "function_id": {"type": "string", "description": "Function id, e.g. investment_portfolio"},
                },
                "required": ["function_id"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "get_function_meta",
            "description": "Get metadata about a function (description, default views, sample metrics).",
            "parameters": {
                "type": "object",
                "properties": {"function_id": {"type": "string"}},
                "required": ["function_id"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "profile_dataset",
            "description": "Statistically profile a dataset: row count, per-column null %, IQR outliers, dtype drift, constant columns.",
            "parameters": {
                "type": "object",
                "properties": {"dataset_id": {"type": "string"}},
                "required": ["dataset_id"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "get_dataset_preview",
            "description": "Return the first N rows of a dataset along with its schema.",
            "parameters": {
                "type": "object",
                "properties": {
                    "dataset_id": {"type": "string"},
                    "n": {"type": "integer", "description": "Max rows (default 25)."},
                },
                "required": ["dataset_id"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "get_model",
            "description": "Get a registered model's full metadata — coefficients, intercept, target, features, train metrics.",
            "parameters": {
                "type": "object",
                "properties": {"model_id": {"type": "string"}},
                "required": ["model_id"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "get_model_metrics",
            "description": "Get a model's monitoring metric trace (R²/AUC/PSI over time).",
            "parameters": {
                "type": "object",
                "properties": {"model_id": {"type": "string"}},
                "required": ["model_id"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "validate_workflow",
            "description": (
                "Validate a workflow design. The analyst's [Context] block contains a `payload` "
                "field that is a JSON string with `nodes` (each {id, kind, ref_id, config}) and "
                "`edges` (each {source, target}). Copy that JSON into `workflow_json` and call. "
                "Returns a list of issues with severity (error / warning / info)."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "workflow_json": {
                        "type": "string",
                        "description": "JSON object with `nodes` and `edges` arrays, copied from the [Context] payload.",
                    },
                },
                "required": ["workflow_json"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "get_run",
            "description": "Fetch a single analytics run by id — status, error, summary, series.",
            "parameters": {
                "type": "object",
                "properties": {"run_id": {"type": "string"}},
                "required": ["run_id"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "preview_tabular_file",
            "description": (
                "Preview a tabular file (.csv / .xlsx / .xls / .parquet) "
                "by absolute path. Returns column names, dtypes, row count, "
                "and the first N rows. Use this when an analyst attached a "
                "data file to a playbook (paths surface in `[UPLOADED FILES]`) "
                "and you need to know the schema before running tools that "
                "consume it (e.g. compute_variance_walk(csv_path=…))."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path":   {"type": "string",  "description": "Absolute path to the file."},
                    "n_rows": {"type": "integer", "description": "Max rows to return (default 10)."},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "get_transform_recipe",
            "description": (
                "Fetch a Transform's full recipe — name, description, inputs, "
                "output_dataset_id, parameters, and the Python source. Use this "
                "before answering 'how does this ETL work' style questions so the "
                "explanation comes from the actual library code, not training data."
            ),
            "parameters": {
                "type": "object",
                "properties": {"transform_id": {"type": "string"}},
                "required": ["transform_id"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "get_tile",
            "description": "Get a tile's saved configuration (plot/table type, source, filters).",
            "parameters": {
                "type": "object",
                "properties": {"tile_id": {"type": "string"}},
                "required": ["tile_id"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "get_tile_preview",
            "description": "Get a tile's live preview rows + column metadata. Use this to know the actual data ranges before suggesting filters.",
            "parameters": {
                "type": "object",
                "properties": {"tile_id": {"type": "string"}},
                "required": ["tile_id"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "apply_tile_filter",
            "description": "Append a filter to a tile so the next preview reflects it. Use this to surface filter chips the analyst can click.",
            "parameters": {
                "type": "object",
                "properties": {
                    "tile_id": {"type": "string"},
                    "field": {"type": "string"},
                    "op": {"type": "string", "enum": ["eq", "ne", "gt", "gte", "lt", "lte", "in", "contains"]},
                    "value": {
                        "description": "Filter value — string, number, boolean, or list of strings/numbers.",
                        "anyOf": [
                            {"type": "string"},
                            {"type": "number"},
                            {"type": "boolean"},
                            {"type": "array", "items": {"anyOf": [{"type": "string"}, {"type": "number"}]}},
                        ],
                    },
                },
                "required": ["tile_id", "field", "op", "value"],
            },
        },
    },
    # ── Plot Tuner mutation tools ─────────────────────────────────────────
    # These work on either a Reporting tile (`target_kind="tile"`,
    # `target_id=<plot_id>`) or an Analytics definition
    # (`target_kind="analytic_def"`, `target_id=<adef_id>`). Each mutates the
    # persisted spec so the next render reflects the change.
    {
        "type": "function",
        "function": {
            "name": "apply_filter",
            "description": "Filter the data underlying a plot or table. Adds one filter row; multiple calls compose AND.",
            "parameters": {
                "type": "object",
                "properties": {
                    "target_kind": {"type": "string", "enum": ["tile", "analytic_def"]},
                    "target_id": {"type": "string"},
                    "field": {"type": "string"},
                    "op": {"type": "string", "enum": ["eq", "ne", "gt", "gte", "lt", "lte", "in", "contains"]},
                    "value": {
                        "description": "Filter value — string, number, boolean, or list of strings/numbers.",
                        "anyOf": [
                            {"type": "string"}, {"type": "number"}, {"type": "boolean"},
                            {"type": "array", "items": {"anyOf": [{"type": "string"}, {"type": "number"}]}},
                        ],
                    },
                },
                "required": ["target_kind", "target_id", "field", "op", "value"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "set_sort",
            "description": "Sort the rendered rows by a field. desc=true is descending. Pass field='' to clear sort.",
            "parameters": {
                "type": "object",
                "properties": {
                    "target_kind": {"type": "string", "enum": ["tile", "analytic_def"]},
                    "target_id": {"type": "string"},
                    "field": {"type": "string"},
                    "desc": {"type": "boolean"},
                },
                "required": ["target_kind", "target_id", "field", "desc"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "set_chart_type",
            "description": "Switch the chart between bar, line, area, stacked_bar, scatter, pie, table, kpi.",
            "parameters": {
                "type": "object",
                "properties": {
                    "target_kind": {"type": "string", "enum": ["tile", "analytic_def"]},
                    "target_id": {"type": "string"},
                    "chart_type": {"type": "string", "enum": ["bar", "line", "area", "stacked_bar", "scatter", "pie", "table", "kpi"]},
                },
                "required": ["target_kind", "target_id", "chart_type"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "set_axes",
            "description": "Pick which field goes on the X axis and which fields are the Y series. Use to add/remove series.",
            "parameters": {
                "type": "object",
                "properties": {
                    "target_kind": {"type": "string", "enum": ["tile", "analytic_def"]},
                    "target_id": {"type": "string"},
                    "x_field": {"type": "string"},
                    "y_fields": {"type": "array", "items": {"type": "string"}},
                },
                "required": ["target_kind", "target_id", "x_field", "y_fields"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "set_axis_labels",
            "description": "Override the chart title and axis labels. Pass empty strings to revert to the auto-derived values.",
            "parameters": {
                "type": "object",
                "properties": {
                    "target_kind": {"type": "string", "enum": ["tile", "analytic_def"]},
                    "target_id": {"type": "string"},
                    "title": {"type": "string"},
                    "x_axis_label": {"type": "string"},
                    "y_axis_label": {"type": "string"},
                },
                "required": ["target_kind", "target_id", "title", "x_axis_label", "y_axis_label"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "set_style",
            "description": "Change visual style — color palette (hex codes in series order), legend position, font size in pixels.",
            "parameters": {
                "type": "object",
                "properties": {
                    "target_kind": {"type": "string", "enum": ["tile", "analytic_def"]},
                    "target_id": {"type": "string"},
                    "palette": {"type": "array", "items": {"type": "string"}, "description": "Hex color list (e.g. ['#0891B2','#7C3AED']). Empty list reverts to default."},
                    "font_size": {"type": "integer", "description": "Pixel size for axis labels + legend. 0 reverts to default."},
                    "legend_position": {"type": "string", "enum": ["top", "bottom", "right", "left", "none"]},
                },
                "required": ["target_kind", "target_id", "palette", "font_size", "legend_position"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "rag_search",
            "description": (
                "Retrieval-Augmented Generation search across a folder of "
                "documents (markdown, text, CSV, XLSX). Returns the top-k "
                "matching chunks ranked by keyword score, each with the "
                "source path and any frontmatter metadata (model_id, "
                "model_component, portfolio_scope) the chunk's parent doc "
                "carries. Use this to surface methodology / whitepaper "
                "context behind a forecast number. doc_dir defaults to "
                "the function's relevant doc folder (configured via the "
                "CMA_DOCS_ROOT env var, fallback `sample_docs/`)."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "query":   {"type": "string",  "description": "Free-text query (multi-word)."},
                    "doc_dir": {"type": "string",  "description": "Optional folder path to scope the search. Empty string = default root, recursive."},
                    "top_k":   {"type": "integer", "description": "Max chunks to return (default 5)."},
                },
                "required": ["query", "doc_dir", "top_k"],
            },
        },
    },
]


# ── Implementations ───────────────────────────────────────────────────────
def _t_get_workspace(args: dict) -> str:
    fid = args.get("function_id", "") or _ctx_function_id()
    ws = get_workspace(fid)
    if not ws:
        return json.dumps({"error": f"No workspace for function `{fid}`"})
    return ws.model_dump_json(indent=2)


def _t_get_function_meta(args: dict) -> str:
    from routers.functions import BUSINESS_FUNCTIONS
    fid = args.get("function_id", "") or _ctx_function_id()
    f = next((x for x in BUSINESS_FUNCTIONS if x.id == fid), None)
    return f.model_dump_json(indent=2) if f else json.dumps({"error": f"Unknown function `{fid}`"})


def _t_profile_dataset(args: dict) -> str:
    did = args.get("dataset_id", "") or _ctx_id_for("dataset")
    d = _DATASETS.get(did)
    if not d:
        return json.dumps({"error": f"Dataset `{did}` not found"})
    df = _read_or_synth(d)
    if df is None or df.empty:
        return json.dumps({"error": f"Could not read dataset `{did}`"})

    profile: dict[str, Any] = {
        "dataset_id": did, "name": d.name, "rows": int(len(df)), "columns": int(len(df.columns)),
        "null_rate": {}, "outliers": {}, "constant_columns": [], "dtype_drift": [],
    }
    for col in df.columns:
        null_pct = float(df[col].isnull().mean() * 100)
        if null_pct > 0:
            profile["null_rate"][col] = round(null_pct, 2)
        if df[col].nunique(dropna=True) <= 1:
            profile["constant_columns"].append(col)

    for col in df.select_dtypes(include="number").columns:
        s = df[col].dropna()
        if len(s) < 4:
            continue
        q1, q3 = float(s.quantile(0.25)), float(s.quantile(0.75))
        iqr = q3 - q1
        if iqr == 0:
            continue
        lo, hi = q1 - 1.5 * iqr, q3 + 1.5 * iqr
        out = s[(s < lo) | (s > hi)]
        if len(out) > 0:
            profile["outliers"][col] = {
                "count": int(len(out)),
                "pct": round(len(out) / len(s) * 100, 2),
                "bounds": [round(lo, 4), round(hi, 4)],
            }

    for col in df.select_dtypes(include="object").columns:
        sample = df[col].dropna().astype(str)
        if sample.empty:
            continue
        as_num = pd.to_numeric(sample, errors="coerce")
        rate = float(as_num.notna().mean())
        if 0.9 <= rate < 1.0:
            profile["dtype_drift"].append({
                "column": col,
                "parse_rate_as_numeric": round(rate * 100, 1),
            })
    return json.dumps(profile, default=_json_default)


def _t_get_dataset_preview(args: dict) -> str:
    did = args.get("dataset_id", "") or _ctx_id_for("dataset")
    n = int(args.get("n", 25))
    d = _DATASETS.get(did)
    if not d:
        # Common LLM mistake: passing a file path / uploaded-file id where a
        # dataset_id is expected. Redirect to preview_tabular_file rather
        # than letting the agent conclude "tools are unavailable".
        looks_like_path = (
            "/" in did or "\\" in did
            or any(did.lower().endswith(ext) for ext in (
                ".csv", ".xlsx", ".xls", ".parquet", ".pdf",
                ".docx", ".pptx", ".md", ".txt", ".json",
            ))
        )
        if looks_like_path:
            return json.dumps({
                "error":           f"`{did}` is a file path, not a dataset id.",
                "fix":             (
                    "Call `preview_tabular_file(path=...)` for files "
                    "(supports CSV / XLSX / Parquet). `get_dataset_preview` "
                    "only works on registered dataset ids that appear in "
                    "the `--- input dataset ---` block of the context."
                ),
                "redirect_tool":   "preview_tabular_file",
            })
        return json.dumps({
            "error":               f"Dataset `{did}` not found",
            "available_datasets":  sorted(list(_DATASETS.keys()))[:20],
            "hint":                (
                "If the analyst attached a file via `[UPLOADED FILES]`, "
                "use `preview_tabular_file(path=…)` instead — uploaded "
                "files are NOT in the dataset registry."
            ),
        })
    df = _read_or_synth(d, n)
    if df is None:
        return json.dumps({"error": "Read failed"})
    return json.dumps({
        "dataset_id": d.id,
        "name": d.name,
        "columns": [{"name": c, "dtype": str(df[c].dtype)} for c in df.columns],
        "rows": _df_records(df, n),
    }, default=_json_default)


def _t_get_model(args: dict) -> str:
    mid = args.get("model_id", "") or _ctx_id_for("model")
    m = _MODELS.get(mid)
    return m.model_dump_json(indent=2) if m else json.dumps({"error": f"Model `{mid}` not found"})


def _t_get_model_metrics(args: dict) -> str:
    mid = args.get("model_id", "") or _ctx_id_for("model")
    m = _MODELS.get(mid)
    if not m:
        return json.dumps({"error": f"Model `{mid}` not found"})
    series: dict[str, list[dict[str, Any]]] = {}
    for entry in m.monitoring_metrics:
        series.setdefault(entry.name, []).append({"asof": entry.asof, "value": entry.value})
    return json.dumps({"model_id": mid, "series": series, "train_metrics": m.train_metrics})


def _t_validate_workflow(args: dict) -> str:
    from routers.chat_validation import validate_workflow_payload  # local helper module
    raw = args.get("workflow_json")
    # Backwards-compat: also accept the older shape with explicit fields
    if raw is None and ("nodes" in args or "edges" in args):
        nodes = args.get("nodes", [])
        edges = args.get("edges", [])
    else:
        if isinstance(raw, str):
            try:
                data = json.loads(raw)
            except json.JSONDecodeError as e:
                return json.dumps({"error": f"workflow_json was not valid JSON: {e}"})
        elif isinstance(raw, dict):
            data = raw
        else:
            return json.dumps({"error": "workflow_json must be a JSON object string"})
        nodes = data.get("nodes", []) or []
        edges = data.get("edges", []) or []
    issues = validate_workflow_payload(nodes, edges)
    return json.dumps({
        "ok": not any(i["severity"] == "error" for i in issues),
        "issues": issues,
        "node_count": len(nodes),
        "edge_count": len(edges),
    })


def _t_get_run(args: dict) -> str:
    rid = args.get("run_id", "") or _ctx_id_for("run")
    r = _RUNS.get(rid)
    return r.model_dump_json(indent=2) if r else json.dumps({"error": f"Run `{rid}` not found"})


def _t_preview_tabular_file(args: dict) -> str:
    """Read a CSV / XLSX / Parquet by absolute path and return shape +
    first N rows. Used when an analyst attached a data file to a
    playbook (paths surface in `[UPLOADED FILES]`) and an agent needs
    to verify the schema before consuming it."""
    import os

    path = (args.get("path") or "").strip()
    if not path:
        return json.dumps({"error": "path is required"})
    n_rows = int(args.get("n_rows") or 10)

    # Allow paths relative to the docs root (so the LLM doesn't have to
    # hand-build the absolute path when only a relative id was surfaced).
    if not os.path.isabs(path):
        try:
            from routers.documents import _docs_root
            cand = (_docs_root() / path).resolve()
            if cand.exists():
                path = str(cand)
        except Exception:
            pass

    if not os.path.exists(path):
        return json.dumps({"error": f"file not found: {path}"})

    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".csv":
            df = pd.read_csv(path)
        elif ext in (".xlsx", ".xls"):
            df = pd.read_excel(path)
        elif ext == ".parquet":
            df = pd.read_parquet(path)
        else:
            return json.dumps({
                "error": f"Unsupported tabular extension {ext!r}",
                "supported": [".csv", ".xlsx", ".xls", ".parquet"],
            })
    except Exception as e:
        return json.dumps({"error": f"read failed: {e}"})

    return json.dumps({
        "path":         path,
        "row_count":    int(len(df)),
        "column_count": int(len(df.columns)),
        "columns":      [{"name": c, "dtype": str(df[c].dtype)} for c in df.columns],
        "sample_rows":  _df_records(df, max(1, min(n_rows, 200))),
        "summary":      _summarize(df),
    }, default=_json_default)


def _t_get_transform_recipe(args: dict) -> str:
    """Return a Transform's full recipe — pulled from the in-process registry
    populated by pack-shipped transforms (`backend/routers/transforms.py`).
    Falls back to the entity-id bound to the current request if the caller
    doesn't pass `transform_id` explicitly."""
    from routers.transforms import _TRANSFORMS

    tid = (args.get("transform_id") or "").strip() or _ctx_id_for("transform")
    if not tid:
        return json.dumps({"error": "transform_id is required (none in args, none on request context)"})
    t = _TRANSFORMS.get(tid)
    if not t:
        return json.dumps({
            "error":              f"Transform `{tid}` not found",
            "available_transforms": sorted(_TRANSFORMS.keys()),
        })
    # Surface the recipe text first so the agent sees real code rather
    # than guessing from the description.
    return json.dumps({
        "id":                    t.id,
        "name":                  t.name,
        "description":           t.description,
        "function_id":           t.function_id,
        "input_data_source_ids": list(t.input_data_source_ids or []),
        "output_dataset_id":     t.output_dataset_id,
        "parameters":            list(t.parameters or []),
        "recipe_python":         t.recipe_python,
        "source":                t.source,
        "pack_id":               t.pack_id,
    }, default=_json_default)


def _t_get_tile(args: dict) -> str:
    tid = args.get("tile_id", "") or _ctx_tile_id()
    p = _PLOTS.get(tid)
    return p.model_dump_json(indent=2) if p else json.dumps({"error": f"Tile `{tid}` not found"})


def _t_get_tile_preview(args: dict) -> str:
    tid = args.get("tile_id", "") or _ctx_tile_id()
    p = _PLOTS.get(tid)
    if not p:
        return json.dumps({"error": f"Tile `{tid}` not found"})
    df = _tile_dataframe(p)
    if df is None or df.empty:
        return json.dumps({"name": p.name, "rows": [], "note": "No live data — using sample"})
    df = _apply_filters(df, p.filters)
    return json.dumps({
        "name": p.name,
        "tile_type": p.tile_type,
        "chart_type": p.chart_type,
        "columns": [{"name": c, "dtype": str(df[c].dtype)} for c in df.columns],
        "row_count": int(len(df)),
        "sample_rows": _df_records(df, 25),
        "summary": _summarize(df),
    }, default=_json_default)


def _t_apply_tile_filter(args: dict) -> str:
    tid = args.get("tile_id", "") or _ctx_tile_id()
    p = _PLOTS.get(tid)
    if not p:
        return json.dumps({"error": f"Tile `{tid}` not found"})
    f = {"field": args.get("field"), "op": args.get("op"), "value": args.get("value")}
    p.filters = list(p.filters) + [f]
    return json.dumps({"ok": True, "filters": p.filters})


# ── Plot-tuner mutation tools (work on tile OR analytic_def) ─────────────
def _resolve_target(args: dict):
    """Return (kind, obj, owner_dict) for whichever registry holds the target.
    Falls back to the active request context's bound entity when the model
    omits target_kind / target_id — so the analyst's "Tune" click is
    sufficient state and the model never has to echo the id back."""
    kind = args.get("target_kind", "") or ""
    tid = args.get("target_id", "") or ""
    if not kind or not tid:
        ck, ce = _ctx_target()
        if ck:
            kind = kind or ck
            tid = tid or ce
    if kind == "tile":
        return ("tile", _PLOTS.get(tid), _PLOTS)
    if kind == "analytic_def":
        from routers.analytics_defs import _DEFS as _ADEFS
        return ("analytic_def", _ADEFS.get(tid), _ADEFS)
    return (kind, None, None)


def _ensure_style(obj):
    """Lazy-init `obj.style` to a PlotStyle if missing — back-compat with
    saved specs from before the field was added."""
    from models.schemas import PlotStyle
    if getattr(obj, "style", None) is None:
        obj.style = PlotStyle()
    return obj.style


def _t_apply_filter(args: dict) -> str:
    kind, obj, _ = _resolve_target(args)
    if obj is None:
        return json.dumps({"error": f"{kind} `{args.get('target_id')}` not found"})
    new_filter = {"field": args.get("field"), "op": args.get("op"), "value": args.get("value")}
    if kind == "tile":
        obj.filters = list(obj.filters) + [new_filter]
        return json.dumps({"ok": True, "kind": kind, "filters": obj.filters})
    if kind == "analytic_def" and obj.kind == "aggregate" and obj.aggregate_spec:
        obj.aggregate_spec.filters = list(obj.aggregate_spec.filters) + [new_filter]
        return json.dumps({"ok": True, "kind": kind, "filters": obj.aggregate_spec.filters})
    return json.dumps({"error": f"Filters not supported on {kind}/{obj.kind}"})


def _t_set_sort(args: dict) -> str:
    kind, obj, _ = _resolve_target(args)
    if obj is None:
        return json.dumps({"error": f"{kind} `{args.get('target_id')}` not found"})
    field = args.get("field", "") or None
    desc = bool(args.get("desc"))
    if kind == "tile":
        # Tiles store table sort separately; apply to style for both plots + tables.
        obj.table_default_sort = field
        obj.table_default_sort_desc = desc
    if kind == "analytic_def" and obj.kind == "aggregate" and obj.aggregate_spec:
        obj.aggregate_spec.sort_by = field
        obj.aggregate_spec.sort_desc = desc
    style = _ensure_style(obj)
    style.sort_field = field
    style.sort_desc = desc
    return json.dumps({"ok": True, "kind": kind, "sort_field": field, "sort_desc": desc})


def _t_set_chart_type(args: dict) -> str:
    kind, obj, _ = _resolve_target(args)
    if obj is None:
        return json.dumps({"error": f"{kind} `{args.get('target_id')}` not found"})
    new_type = args.get("chart_type", "")
    if kind == "tile":
        # PlotConfig restricts chart_type to a smaller enum; KPI/table maps to tile_type
        if new_type in ("table",):
            obj.tile_type = "table"
        else:
            obj.tile_type = "plot"
            obj.chart_type = new_type if new_type in (
                "line", "bar", "area", "pie", "scatter", "stacked_bar"
            ) else obj.chart_type
        return json.dumps({"ok": True, "kind": kind, "chart_type": obj.chart_type, "tile_type": obj.tile_type})
    if kind == "analytic_def" and obj.output:
        obj.output.chart_type = new_type
        return json.dumps({"ok": True, "kind": kind, "chart_type": obj.output.chart_type})
    return json.dumps({"error": f"Cannot set chart_type on {kind}"})


def _t_set_axes(args: dict) -> str:
    kind, obj, _ = _resolve_target(args)
    if obj is None:
        return json.dumps({"error": f"{kind} `{args.get('target_id')}` not found"})
    x = args.get("x_field", "")
    ys = args.get("y_fields", []) or []
    if kind == "tile":
        obj.x_field = x or obj.x_field
        obj.y_fields = list(ys) if ys else obj.y_fields
        return json.dumps({"ok": True, "kind": kind, "x_field": obj.x_field, "y_fields": obj.y_fields})
    if kind == "analytic_def" and obj.output:
        obj.output.x_field = x or obj.output.x_field
        obj.output.y_fields = list(ys) if ys else obj.output.y_fields
        return json.dumps({"ok": True, "kind": kind, "x_field": obj.output.x_field, "y_fields": obj.output.y_fields})
    return json.dumps({"error": f"Cannot set axes on {kind}"})


def _t_set_axis_labels(args: dict) -> str:
    kind, obj, _ = _resolve_target(args)
    if obj is None:
        return json.dumps({"error": f"{kind} `{args.get('target_id')}` not found"})
    style = _ensure_style(obj)
    title = args.get("title")
    xl = args.get("x_axis_label")
    yl = args.get("y_axis_label")
    # Empty string explicitly clears (reverts to auto-derived).
    style.title = title if title else None
    style.x_axis_label = xl if xl else None
    style.y_axis_label = yl if yl else None
    return json.dumps({"ok": True, "kind": kind, "title": style.title,
                       "x_axis_label": style.x_axis_label, "y_axis_label": style.y_axis_label})


def _t_set_style(args: dict) -> str:
    kind, obj, _ = _resolve_target(args)
    if obj is None:
        return json.dumps({"error": f"{kind} `{args.get('target_id')}` not found"})
    style = _ensure_style(obj)
    palette = args.get("palette") or []
    font_size = args.get("font_size") or 0
    legend = args.get("legend_position", "")
    style.palette = list(palette) if palette else []
    style.font_size = int(font_size) if font_size and font_size > 0 else None
    if legend in ("top", "bottom", "right", "left", "none"):
        style.legend_position = legend
    return json.dumps({"ok": True, "kind": kind,
                       "palette": style.palette,
                       "font_size": style.font_size,
                       "legend_position": style.legend_position})


# ── RAG search (built-in) ─────────────────────────────────────────────────
# Generic keyword-scoring RAG over a folder of documents. Replaces the
# pack-specific `search_methodology_docs` tool — any agent can use this
# without each pack re-implementing the same loop. Demo-grade scoring;
# swap the body for a real vector store in production.
def _t_rag_search(args: dict) -> str:
    import glob
    import os
    import re

    query = (args.get("query") or "").strip()
    if not query:
        return json.dumps({"matches": [], "note": "empty query"})

    doc_dir = (args.get("doc_dir") or "").strip()
    if not doc_dir:
        # Default to the configured docs root, then recurse.
        env_root = os.getenv("CMA_DOCS_ROOT", "").strip()
        if env_root and os.path.isdir(env_root):
            doc_dir = env_root
        else:
            # Repo-relative default: <repo>/sample_docs/
            here = os.path.dirname(os.path.abspath(__file__))
            repo = os.path.abspath(os.path.join(here, "..", ".."))
            doc_dir = os.path.join(repo, "sample_docs")

    if not os.path.isdir(doc_dir):
        return json.dumps({"error": f"doc_dir does not exist: {doc_dir}"})

    q_tokens = [t for t in re.findall(r"[a-z0-9_]+", query.lower()) if len(t) > 2]
    if not q_tokens:
        return json.dumps({"matches": []})

    def _read_text(path: str) -> str:
        """Best-effort text extraction across the formats the Knowledge Base
        accepts. Optional libraries (`pypdf`, `python-docx`, `python-pptx`)
        are imported lazily so a missing library degrades to "skip this
        file" rather than crashing the whole search."""
        ext = os.path.splitext(path)[1].lower()
        try:
            if ext in (".md", ".txt", ".py", ".json"):
                with open(path, encoding="utf-8", errors="ignore") as fh:
                    return fh.read()
            if ext == ".csv":
                return pd.read_csv(path).to_csv(index=False)
            if ext in (".xlsx", ".xls"):
                return pd.read_excel(path).to_csv(index=False)
            if ext == ".pdf":
                try:
                    from pypdf import PdfReader
                except ImportError:
                    return ""
                reader = PdfReader(path)
                return "\n\n".join((p.extract_text() or "") for p in reader.pages)
            if ext == ".docx":
                try:
                    from docx import Document  # python-docx
                except ImportError:
                    return ""
                doc = Document(path)
                paras = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
                # Pull table cell text too — whitepapers love tables.
                for tbl in doc.tables:
                    for row in tbl.rows:
                        for cell in row.cells:
                            if cell.text and cell.text.strip():
                                paras.append(cell.text.strip())
                return "\n\n".join(paras)
            if ext == ".pptx":
                try:
                    from pptx import Presentation  # python-pptx
                except ImportError:
                    return ""
                prs = Presentation(path)
                slides_out: list[str] = []
                for idx, slide in enumerate(prs.slides, start=1):
                    parts = [f"[Slide {idx}]"]
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                txt = "".join(run.text for run in para.runs)
                                if txt and txt.strip():
                                    parts.append(txt.strip())
                    if len(parts) > 1:
                        slides_out.append("\n".join(parts))
                # Each slide becomes its own paragraph-block so the chunker
                # treats it as a single semantic unit.
                return "\n\n".join(slides_out)
        except Exception:
            return ""
        return ""

    def _parse_md(text: str) -> tuple[dict[str, str], list[str]]:
        meta: dict[str, str] = {}
        m = re.match(r"^---\n(.*?)\n---\n(.*)$", text, re.DOTALL)
        body = m.group(2) if m else text
        if m:
            for line in m.group(1).splitlines():
                if ":" in line:
                    k, _, v = line.partition(":")
                    meta[k.strip()] = v.strip().strip('"')
        chunks = [c.strip() for c in re.split(r"\n\s*\n", body) if c.strip()]
        return meta, chunks

    top_k = int(args.get("top_k") or 5)

    scored: list[dict[str, Any]] = []
    patterns = (
        "*.md", "*.txt", "*.py", "*.json",
        "*.csv", "*.xlsx", "*.xls",
        "*.pdf", "*.docx", "*.pptx",
    )
    seen_paths: set[str] = set()
    for pat in patterns:
        for path in sorted(glob.glob(os.path.join(doc_dir, "**", pat), recursive=True)):
            if path in seen_paths:
                continue
            seen_paths.add(path)
            text = _read_text(path)
            if not text:
                continue
            if path.endswith(".md"):
                meta, chunks = _parse_md(text)
            else:
                # Split non-md text on blank lines so PDF pages, DOCX
                # paragraphs, and PPTX slides each become their own chunk
                # (read_text emits `\n\n` between those boundaries).
                meta, chunks = {}, [c.strip() for c in re.split(r"\n\s*\n", text) if c.strip()]
            for i, chunk in enumerate(chunks):
                tokens = re.findall(r"[a-z0-9_]+", chunk.lower())
                score = sum(tokens.count(t) for t in q_tokens)
                if score == 0:
                    continue
                preview = chunk[:480] + ("…" if len(chunk) > 480 else "")
                scored.append({
                    "score":           score,
                    "model_id":        meta.get("model_id", os.path.basename(path)),
                    "model_component": meta.get("model_component", ""),
                    "portfolio_scope": meta.get("portfolio_scope", ""),
                    "doc_path":        path,
                    "chunk_index":     i,
                    "preview":         preview,
                })
    scored.sort(key=lambda r: -r["score"])
    return json.dumps({"query": query, "doc_dir": doc_dir, "matches": scored[:top_k]})


# ── Handler ───────────────────────────────────────────────────────────────
_HANDLERS = {
    "get_workspace":      _t_get_workspace,
    "get_function_meta":  _t_get_function_meta,
    "profile_dataset":    _t_profile_dataset,
    "get_dataset_preview": _t_get_dataset_preview,
    "get_model":          _t_get_model,
    "get_model_metrics":  _t_get_model_metrics,
    "validate_workflow":  _t_validate_workflow,
    "get_run":            _t_get_run,
    "get_transform_recipe": _t_get_transform_recipe,
    "preview_tabular_file": _t_preview_tabular_file,
    "get_tile":           _t_get_tile,
    "get_tile_preview":   _t_get_tile_preview,
    "apply_tile_filter":  _t_apply_tile_filter,
    # Plot-tuner mutation tools — work on tiles OR analytic-def specs
    "apply_filter":       _t_apply_filter,
    "set_sort":           _t_set_sort,
    "set_chart_type":     _t_set_chart_type,
    "set_axes":           _t_set_axes,
    "set_axis_labels":    _t_set_axis_labels,
    "set_style":          _t_set_style,
    # Generic RAG retrieval — replaces the deposits-pack scoped variant.
    "rag_search":         _t_rag_search,
}


def handle_tool_call(name: str, args: dict) -> str:
    handler = _HANDLERS.get(name)
    if not handler:
        return json.dumps({"error": f"Unknown tool `{name}`"})
    try:
        return handler(args)
    except Exception as e:
        import traceback
        return json.dumps({"error": str(e), "traceback": traceback.format_exc()[-500:]})


# ── Helpers ───────────────────────────────────────────────────────────────
def _read_or_synth(d, n: int = 2000) -> pd.DataFrame | None:
    if d.source_kind == "upload" and d.file_path and d.file_format:
        try:
            return _read_dataframe(_resolve_path(d), d.file_format).head(n)
        except Exception:
            return None
    rows = _synthesize_sample(d.columns, n)
    return pd.DataFrame(rows)


def _tile_dataframe(p) -> pd.DataFrame | None:
    if p.dataset_id:
        d = _DATASETS.get(p.dataset_id)
        if d and d.source_kind == "upload" and d.file_path and d.file_format:
            try:
                return _read_dataframe(_resolve_path(d), d.file_format)
            except Exception:
                return None
    if p.run_id:
        run = _RUNS.get(p.run_id)
        if run and run.series:
            return pd.DataFrame(run.series)
    return None


def _df_records(df: pd.DataFrame, n: int) -> list[dict[str, Any]]:
    out = df.head(n).copy()
    for c in out.columns:
        if pd.api.types.is_datetime64_any_dtype(out[c]):
            out[c] = out[c].astype(str)
    return json.loads(out.to_json(orient="records", date_format="iso"))


def _summarize(df: pd.DataFrame) -> dict[str, Any]:
    summary: dict[str, Any] = {}
    for col in df.select_dtypes(include="number").columns[:6]:
        s = df[col].dropna()
        if len(s) == 0:
            continue
        summary[col] = {
            "min": float(s.min()), "max": float(s.max()),
            "median": float(s.median()), "p90": float(s.quantile(0.9)),
        }
    for col in df.select_dtypes(include="object").columns[:3]:
        vc = df[col].value_counts().head(3)
        summary[col] = {str(k): int(v) for k, v in vc.items()}
    return summary


def _json_default(o):
    if isinstance(o, float) and (math.isnan(o) or math.isinf(o)):
        return None
    return str(o)
